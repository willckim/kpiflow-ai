# shared/analysis.py
import os
import io
from typing import Dict, List, Any, Tuple, Optional

import numpy as np
import pandas as pd

# --- .env support (optional) ---
try:
    from dotenv import load_dotenv  # type: ignore
    load_dotenv()
except Exception:
    pass

# --- OpenAI / Azure OpenAI availability flags (optional) ---
OPENAI_AVAILABLE = False
AZURE_OPENAI_AVAILABLE = False
try:
    from openai import OpenAI  # direct OpenAI
    OPENAI_AVAILABLE = True
except Exception:
    OPENAI_AVAILABLE = False

try:
    from openai import AzureOpenAI  # Azure OpenAI client (same package)
    AZURE_OPENAI_AVAILABLE = True
except Exception:
    AZURE_OPENAI_AVAILABLE = False


# =========================
# Data loading & cleaning
# =========================
def load_data(file_bytes: bytes, filename: str) -> pd.DataFrame:
    """Load CSV or Excel from raw bytes."""
    name = filename.lower()
    if name.endswith(".csv"):
        return pd.read_csv(io.BytesIO(file_bytes))
    if name.endswith((".xls", ".xlsx")):
        return pd.read_excel(io.BytesIO(file_bytes))
    raise ValueError("Unsupported file type. Please upload CSV or Excel.")


def clean_data(df: pd.DataFrame) -> pd.DataFrame:
    """Light normalization: trim/lower column names, coerce dates, fill NA."""
    df = df.copy()

    # Normalize column names
    df.columns = [c.strip() for c in df.columns]
    lower_map = {c.lower(): c for c in df.columns}
    # Rebuild columns preserving original case but expose lowercase mapping
    # (we'll use lowercase lookups but keep original data)
    df.columns = [lower_map.get(c.lower(), c) for c in df.columns]

    # Try to coerce any column with 'date' in name to datetime
    for c in df.columns:
        if "date" in c.lower():
            try:
                df[c] = pd.to_datetime(df[c], errors="coerce")
            except Exception:
                pass

    # Fill NA: numeric → 0, object → "Unknown"
    for c in df.select_dtypes(include=[np.number]).columns:
        df[c] = df[c].fillna(0)
    for c in df.select_dtypes(include=["object"]).columns:
        df[c] = df[c].fillna("Unknown")

    return df


# =========================
# KPI computation helpers
# =========================
def _maybe_col(df: pd.DataFrame, candidates: List[str]) -> str:
    """Return the first matching column name (case-insensitive) or ''."""
    cols_lower = {c.lower(): c for c in df.columns}
    for cand in candidates:
        if cand in cols_lower:
            return cols_lower[cand]
    return ""


def _monthly_series(df: pd.DataFrame, date_col: str, value_col: str) -> pd.Series:
    s = df.dropna(subset=[date_col]).copy()
    if s.empty:
        return pd.Series(dtype=float)
    s["__month"] = pd.to_datetime(s[date_col], errors="coerce").dt.to_period("M")
    monthly = s.groupby("__month")[value_col].sum().sort_index()
    # Ensure float so we can cast cleanly later
    return monthly.astype(float)


# =========================
# KPIs (match Streamlit keys)
# =========================
def compute_kpis(df: pd.DataFrame) -> Dict[str, Any]:
    """
    Returns keys aligned with the Streamlit frontend:
      revenue, gross_profit, gross_margin_pct, unique_customers,
      mom_growth_pct, last_month_revenue, avg_order_value, orders
    """
    revenue_col = _maybe_col(df, ["revenue", "sales", "amount", "total", "price"])
    cost_col    = _maybe_col(df, ["cost", "cogs"])
    date_col    = _maybe_col(df, ["order_date", "date", "created_at"])
    cust_col    = _maybe_col(df, ["customer_id", "customer", "user_id"])
    order_col   = _maybe_col(df, ["order_id", "order", "invoice_id", "transaction_id"])

    out: Dict[str, Any] = {
        "revenue": None,
        "gross_profit": None,
        "gross_margin_pct": None,
        "unique_customers": None,
        "mom_growth_pct": None,
        "last_month_revenue": None,
        "avg_order_value": None,
        "orders": None,
    }

    # Revenue
    if revenue_col:
        total_rev = float(df[revenue_col].sum())
        out["revenue"] = total_rev
    else:
        return out  # cannot compute anything meaningful without revenue

    # Gross profit & margin
    if cost_col:
        gross_profit = float(df[revenue_col].sum() - df[cost_col].sum())
        out["gross_profit"] = gross_profit
        out["gross_margin_pct"] = float(100.0 * gross_profit / df[revenue_col].sum()) if df[revenue_col].sum() else None

    # Unique customers
    if cust_col:
        out["unique_customers"] = int(df[cust_col].nunique())

    # Orders & AOV
    if order_col:
        orders = int(df[order_col].nunique())
    else:
        # fallback: treat each row as an order if no explicit order id
        orders = int(len(df))
    out["orders"] = orders
    out["avg_order_value"] = float(out["revenue"] / orders) if orders > 0 else None

    # Month-over-month & last month revenue
    if date_col:
        monthly = _monthly_series(df, date_col, revenue_col)
        if len(monthly) >= 1:
            out["last_month_revenue"] = float(monthly.iloc[-1])
        if len(monthly) >= 2 and monthly.iloc[-2] != 0:
            out["mom_growth_pct"] = float(100.0 * (monthly.iloc[-1] - monthly.iloc[-2]) / monthly.iloc[-2])

    return out


# =========================
# Optional charts (lazy imports so backend can stay lean)
# =========================
def build_charts(df: pd.DataFrame) -> Dict[str, Any]:
    """
    Returns Plotly figures. Lazy-imports plotly so this module can be used
    without plotly installed (unless you actually call this function).
    """
    import plotly.express as px  # lazy import

    charts: Dict[str, Any] = {}
    revenue_col = _maybe_col(df, ["revenue", "sales", "amount", "total", "price"])
    cost_col    = _maybe_col(df, ["cost", "cogs"])
    date_col    = _maybe_col(df, ["order_date", "date", "created_at"])
    region_col  = _maybe_col(df, ["region", "state", "country"])
    category_col= _maybe_col(df, ["category", "product", "item"])

    if date_col and revenue_col:
        d = df.dropna(subset=[date_col]).copy()
        if not d.empty:
            d["__date"] = pd.to_datetime(d[date_col], errors="coerce")
            ts = d.groupby(pd.Grouper(key="__date", freq="W"))[revenue_col].sum().reset_index()
            charts["revenue_trend"] = px.line(ts, x="__date", y=revenue_col, title="Revenue Trend (Weekly)")

    if region_col and revenue_col:
        region = df.groupby(region_col, dropna=False)[revenue_col].sum().reset_index().sort_values(revenue_col, ascending=False)
        charts["by_region"] = px.bar(region, x=region_col, y=revenue_col, title="Revenue by Region")

    if category_col and revenue_col:
        cat = df.groupby(category_col, dropna=False)[revenue_col].sum().reset_index().sort_values(revenue_col, ascending=False)
        charts["by_category"] = px.bar(cat, x=category_col, y=revenue_col, title="Revenue by Category")

    if revenue_col and cost_col:
        df2 = df.copy()
        df2["__gross_profit"] = df2[revenue_col] - df2[cost_col]
        agg = df2[[revenue_col, "__gross_profit"]].sum()
        charts["profit_bars"] = px.bar(
            x=["Revenue", "Gross Profit"],
            y=[agg[revenue_col], agg["__gross_profit"]],
            title="Revenue vs Gross Profit",
        )
    return charts


# =========================
# Insights (LLM optional)
# =========================
def _format_kpis_for_prompt(kpis: Dict[str, Any]) -> str:
    label_map = {
        "revenue": "Revenue",
        "gross_profit": "Gross Profit",
        "gross_margin_pct": "Gross Margin %",
        "unique_customers": "Unique Customers",
        "mom_growth_pct": "MoM Growth %",
        "last_month_revenue": "Last Month Revenue",
        "avg_order_value": "Avg Order Value",
        "orders": "Orders",
    }
    parts = []
    for k, v in kpis.items():
        label = label_map.get(k, k)
        if isinstance(v, float):
            parts.append(f"- {label}: {round(v, 2)}")
        else:
            parts.append(f"- {label}: {v}")
    return "\n".join(parts)


def _ai_client_and_model() -> Tuple[Optional[Any], Optional[str], bool]:
    """Return (client, model_name, is_azure). Prefers Azure if fully configured."""
    # Azure first
    az_key   = os.getenv("AZURE_OPENAI_KEY")
    az_ep    = os.getenv("AZURE_OPENAI_ENDPOINT")
    az_ver   = os.getenv("AZURE_OPENAI_API_VERSION")
    az_deploy= os.getenv("AZURE_OPENAI_DEPLOYMENT")

    if AZURE_OPENAI_AVAILABLE and az_key and az_ep and az_ver and az_deploy:
        try:
            client = AzureOpenAI(api_key=az_key, api_version=az_ver, azure_endpoint=az_ep)
            return client, az_deploy, True
        except Exception:
            pass

    # Direct OpenAI
    openai_key = os.getenv("OPENAI_API_KEY")
    if OPENAI_AVAILABLE and openai_key:
        try:
            client = OpenAI(api_key=openai_key)
            model = os.getenv("OPENAI_MODEL", "gpt-5-mini")
            return client, model, False
        except Exception:
            pass

    return None, None, False


def generate_insights(df: pd.DataFrame, kpis: Dict[str, Any]) -> str:
    """Try LLM for insights; otherwise provide heuristic bullets."""
    client, model, is_azure = _ai_client_and_model()
    if client is None or model is None:
        return _heuristic_insights(kpis)

    prompt = f"""You are an analytics copilot. Given the KPIs below, write 3–6 executive-ready bullet insights.
KPIs:
{_format_kpis_for_prompt(kpis)}
Keep it concise, factual, and helpful. Avoid overclaiming causality."""
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You write precise business insights."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.2,
            max_tokens=250,
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception as e:
        return _heuristic_insights(kpis, note=f"AI insights unavailable ({type(e).__name__}).")


def _heuristic_insights(kpis: Dict[str, Any], note: Optional[str] = None) -> str:
    """Human-readable fallback if LLM is not configured or errors out."""
    lines: List[str] = []
    if note:
        lines.append(note)

    g = kpis.get("mom_growth_pct")
    if isinstance(g, (int, float)):
        if g > 0:
            lines.append(f"Revenue increased {g:.1f}% month-over-month, indicating positive momentum.")
        elif g < 0:
            lines.append(f"Revenue decreased {abs(g):.1f}% month-over-month; investigate potential causes.")

    gm = kpis.get("gross_margin_pct")
    if isinstance(gm, (int, float)):
        lines.append(f"Gross margin is {gm:.1f}%. Consider pricing or cost optimization to improve profitability.")

    if not lines:
        lines.append("Data processed successfully. Add Azure/OpenAI keys to unlock richer AI insights.")
    return "\n".join(lines)


# =========================
# Optional PPTX export (lazy imports)
# =========================
def export_pptx(title: str, kpis: Dict[str, Any], insights: str, charts: Dict[str, Any], out_path: str):
    """
    Create a PPTX with KPI bullets and Plotly charts.
    Lazy-imports python-pptx and kaleido/plotly.io so backend can remain lean.
    """
    from pptx import Presentation  # lazy import
    from pptx.util import Inches
    import plotly.io as pio        # lazy import
    import os as _os

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "KPIFlow AI — Auto-generated KPI & Insights"

    # KPI slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPI Summary"
    body = slide.shapes.placeholders[1].text_frame
    for k, v in kpis.items():
        body.add_paragraph().text = f"{k}: {round(v, 2) if isinstance(v, float) else v}"

    # Insights slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Insights"
    body = slide.shapes.placeholders[1].text_frame
    for line in (insights or "").split("\n"):
        if line.strip():
            body.add_paragraph().text = line.strip()

    # Chart slides
    tmp_dir = _os.path.join(_os.path.dirname(out_path) or ".", "_tmp_imgs")
    _os.makedirs(tmp_dir, exist_ok=True)
    for name, fig in (charts or {}).items():
        img_path = _os.path.join(tmp_dir, f"{name}.png")
        pio.write_image(fig, img_path, format="png", width=1280, height=720, scale=2)

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = name.replace("_", " ").title()
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))

    prs.save(out_path)
